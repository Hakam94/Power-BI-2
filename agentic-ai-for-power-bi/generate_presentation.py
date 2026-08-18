"""
Hakam Data Studio — Agentic AI for Power BI
PowerPoint Generator v2 — FULL REDESIGN
Proper split layouts, no overlapping, proper visual storytelling
Brand: Cyber Slate #0D2229 | Lime Green #BFFF00 | Electric Cyan #00D2FF
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ─── Brand Colors ──────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x22, 0x29)  # Cyber Slate
CARD    = RGBColor(0x0A, 0x18, 0x1F)  # Darker card
LIME    = RGBColor(0xBF, 0xFF, 0x00)  # Lime Green
CYAN    = RGBColor(0x00, 0xD2, 0xFF)  # Electric Cyan
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DIMMED  = RGBColor(0x88, 0xAA, 0xBB)
RED_DIM = RGBColor(0x66, 0x11, 0x11)
GRN_DIM = RGBColor(0x0A, 0x28, 0x10)

W = Inches(13.33)
H = Inches(7.5)

SCRIPT_DIR = Path(__file__).parent

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, line_color=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, qn('a:pPr'))
        pPr.set('rtl', '1' if rtl else '0')
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name  = "Arial"
    return box

def img(slide, filename, l, t, w, h=None):
    p = SCRIPT_DIR / filename
    if p.exists():
        if h:
            slide.shapes.add_picture(str(p), l, t, w, h)
        else:
            slide.shapes.add_picture(str(p), l, t, w)
    else:
        # placeholder gray box
        rect(slide, l, t, w, h or Inches(3), DIMMED)

def bar(slide, top, color=LIME, height=Inches(0.06)):
    rect(slide, Inches(0), top, W, height, color)

def logo(slide):
    txt(slide, "@HakamDataStudio", Inches(0.3), Inches(7.05),
        Inches(3), Inches(0.35), size=11, bold=True, color=CYAN,
        align=PP_ALIGN.LEFT, rtl=False)

def num(slide, n):
    txt(slide, str(n), Inches(12.7), Inches(7.05),
        Inches(0.5), Inches(0.35), size=11, color=DIMMED,
        align=PP_ALIGN.CENTER, rtl=False)


# ══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), CYAN)
bar(s, Inches(7.44), LIME)

# Left accent stripe
rect(s, Inches(0), Inches(0.06), Inches(0.08), Inches(7.38), LIME)

# Right dark panel
rect(s, Inches(6.8), Inches(0.06), Inches(6.53), Inches(7.38), CARD)

# Left side — channel branding
txt(s, "Hakam Data Studio", Inches(0.3), Inches(1.2), Inches(6.2), Inches(0.6),
    size=20, color=CYAN, align=PP_ALIGN.LEFT, rtl=False)
txt(s, "@HakamDataStudio", Inches(0.3), Inches(1.8), Inches(6.2), Inches(0.5),
    size=14, color=DIMMED, align=PP_ALIGN.LEFT, rtl=False)
rect(s, Inches(0.3), Inches(2.4), Inches(5.5), Inches(0.04), LIME)
txt(s, "مؤتمر مايكروسوفت بيلد 2026\n(Microsoft Build 2026)",
    Inches(0.3), Inches(2.5), Inches(6.2), Inches(0.9),
    size=16, color=DIMMED, align=PP_ALIGN.LEFT, rtl=False)
txt(s, "شرح كامل بالعربي", Inches(0.3), Inches(6.5), Inches(6.2), Inches(0.5),
    size=16, color=LIME, align=PP_ALIGN.LEFT, rtl=False)

# Right side — main title
txt(s, "عصر الذكاء\nالاصطناعي الوكيل",
    Inches(7.0), Inches(1.0), Inches(6.0), Inches(2.8),
    size=46, bold=True, color=WHITE)
txt(s, "(Agentic AI Era)",
    Inches(7.0), Inches(3.7), Inches(6.0), Inches(0.7),
    size=24, bold=True, color=LIME)
rect(s, Inches(7.0), Inches(4.5), Inches(5.8), Inches(0.04), CYAN)
txt(s, "في باور بي آي ومايكروسوفت فابريك",
    Inches(7.0), Inches(4.65), Inches(6.0), Inches(0.65),
    size=18, color=CYAN)
txt(s, "(Power BI & Microsoft Fabric)",
    Inches(7.0), Inches(5.3), Inches(6.0), Inches(0.5),
    size=14, color=DIMMED)

logo(s); num(s, 1)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — 3 ERAS: full image slide with overlay cards
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# ═══════════════════════════════════════════════════
# SLIDE 2 — 3 ERAS: Left image panel | Right cards
# ═══════════════════════════════════════════════════
# LEFT PANEL — dark bg, image centered with correct aspect ratio
rect(s, Inches(0), Inches(0.06), Inches(6.6), Inches(7.38), CARD)
bar(s, Inches(0), LIME)

txt(s, "التطور عبر ثلاث مراحل",
    Inches(0.2), Inches(0.15), Inches(6.2), Inches(0.55),
    size=20, bold=True, color=LIME, align=PP_ALIGN.RIGHT)
txt(s, "(Three Eras of Analytics)",
    Inches(0.2), Inches(0.68), Inches(6.2), Inches(0.38),
    size=13, color=CYAN, align=PP_ALIGN.RIGHT, rtl=False)

# Image: 6" wide, auto height to preserve aspect ratio (likely ~6" tall for square)
# Place it centered vertically in the left panel
img(s, "agentic_stack_diagram.png",
    Inches(0.3), Inches(1.1), Inches(6.0))   # width only — auto height

# RIGHT PANEL — title + 3 stacked era cards
rect(s, Inches(6.7), Inches(0.06), Inches(6.63), Inches(7.38), BG)
txt(s, "من لوحات البيانات",
    Inches(6.9), Inches(0.15), Inches(6.3), Inches(0.55),
    size=22, bold=True, color=WHITE)
txt(s, "إلى الوكلاء الذكيين المستقلين",
    Inches(6.9), Inches(0.68), Inches(6.3), Inches(0.5),
    size=18, color=LIME)
rect(s, Inches(6.8), Inches(1.25), Inches(6.2), Inches(0.04), LIME)

eras = [
    ("+20 سنة", "المرحلة الأولى  (Dashboard Era)",
     "لوحات ثابتة — الإنسان يسأل والنظام يجاوب\nنموذج سلبي في اتجاه واحد",
     LIME, "📊"),
    ("2022-2025", "المرحلة الثانية  (Copilot Era)",
     "مساعد ذكي — يجاوب الأسئلة بالكلام الطبيعي\nلكن الإنسان يقرر وينفذ",
     CYAN, "🤝"),
    ("2026 ←", "المرحلة الثالثة  (Agentic AI Era)",
     "وكيل مستقل — يخطط وينفذ ويحقق الأهداف\nبدون تدخل بشري في كل خطوة",
     LIME, "🚀"),
]
for i, (yr, title, desc, col, icon) in enumerate(eras):
    top = Inches(1.38 + i * 2.0)
    rect(s, Inches(6.8), top, Inches(6.2), Inches(1.82), CARD)
    rect(s, Inches(6.8), top, Inches(0.06), Inches(1.82), col)
    # Icon + Year badge
    txt(s, icon, Inches(6.9), top + Inches(0.12),
        Inches(0.55), Inches(0.55), size=22, align=PP_ALIGN.LEFT, rtl=False)
    txt(s, yr, Inches(7.5), top + Inches(0.15), Inches(5.2), Inches(0.42),
        size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)
    # Title
    txt(s, title, Inches(6.9), top + Inches(0.62), Inches(5.9), Inches(0.5),
        size=14, bold=True, color=WHITE)
    # Desc
    txt(s, desc, Inches(6.9), top + Inches(1.1), Inches(5.9), Inches(0.65),
        size=11, color=DIMMED)

logo(s); num(s, 2)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — 3 BIG CHANGES (icon-driven cards, no image overlap)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), CYAN)
bar(s, Inches(7.44), LIME)

txt(s, "مؤتمر مايكروسوفت بيلد 2026 — ثلاثة تغييرات كبرى",
    Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.65),
    size=26, bold=True, color=LIME)
txt(s, "(Microsoft Build 2026 — 3 Major Announcements)",
    Inches(0.4), Inches(0.78), Inches(12.5), Inches(0.38),
    size=15, color=CYAN, align=PP_ALIGN.RIGHT, rtl=False)

changes = [
    ("01", "🤖",
     "مهارات الوكيل\n(Agent Skills for Power BI)",
     "بناء موديل أعمال كامل وتقارير\nبالكلام الطبيعي أو لقطة شاشة\n(Screenshot)",
     ["✦ موديل الأعمال  (Semantic Model)",
      "✦ مقاييس الداكس  (DAX Measures)",
      "✦ صفحات التقرير  (Report Pages)",
      "✦ نشر في الفابريك  (Publish to Fabric)"],
     LIME),
    ("02", "🌐",
     "تطبيقات فابريك\n(Fabric Apps for Semantic Models)",
     "تطبيقات ويب أصيلة فوق موديل\nالأعمال — بدون فريق تطوير\n(No Dev Team Needed)",
     ["✦ واجهة أمامية كاملة  (Full Frontend)",
      "✦ نشر مباشر  (Direct Publish)",
      "✦ مبني على الموديل  (Model-Driven)",
      "✦ في أيام وليس أسابيع"],
     CYAN),
    ("03", "🛒",
     "متجر المهارات\n(Skills for Fabric Marketplace)",
     "سوق متكامل لأدوات الوكلاء\nالذكيين مع حوكمة كاملة\n(Fully Governed)",
     ["✦ حزم جاهزة  (Ready Packages)",
      "✦ متوافق مع كل الوكلاء",
      "✦ حوكمة مؤسسية  (Enterprise Gov.)",
      "✦ أمان كامل  (Full Security)"],
     LIME),
]

for i, (num_lbl, icon, title, desc, bullets, col) in enumerate(changes):
    l = Inches(0.35 + i * 4.35)
    # Card background
    rect(s, l, Inches(1.3), Inches(4.2), Inches(5.95), CARD)
    rect(s, l, Inches(1.3), Inches(4.2), Inches(0.07), col)

    # Number badge
    rect(s, l + Inches(0.15), Inches(1.4), Inches(0.55), Inches(0.55), col)
    txt(s, num_lbl, l + Inches(0.15), Inches(1.4), Inches(0.55), Inches(0.55),
        size=16, bold=True, color=BG, align=PP_ALIGN.CENTER, rtl=False)

    # Icon
    txt(s, icon, l, Inches(1.42), Inches(4.2), Inches(0.7),
        size=30, align=PP_ALIGN.CENTER, rtl=False)

    # Title
    txt(s, title, l, Inches(2.05), Inches(4.2), Inches(1.0),
        size=16, bold=True, color=col, align=PP_ALIGN.CENTER, rtl=False)

    # Divider
    rect(s, l + Inches(0.2), Inches(3.05), Inches(3.8), Inches(0.03), col)

    # Description
    txt(s, desc, l, Inches(3.1), Inches(4.2), Inches(1.1),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # Bullet points
    for j, b in enumerate(bullets):
        txt(s, b, l + Inches(0.2), Inches(4.22 + j * 0.58),
            Inches(3.9), Inches(0.5),
            size=12, color=DIMMED)

logo(s); num(s, 3)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — AGENT SKILLS: LEFT IMAGE | RIGHT TEXT
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), LIME)

# Left half — diagram (6.5 wide)
rect(s, Inches(0), Inches(0.06), Inches(6.5), Inches(7.38), CARD)
txt(s, "تدفق عمل مهارات الوكيل",
    Inches(0.2), Inches(0.15), Inches(6.1), Inches(0.55),
    size=18, bold=True, color=LIME, align=PP_ALIGN.RIGHT)
txt(s, "(Agent Skills Workflow)",
    Inches(0.2), Inches(0.68), Inches(6.1), Inches(0.38),
    size=13, color=CYAN, align=PP_ALIGN.RIGHT, rtl=False)
img(s, "agent_skills_workflow.png",
    Inches(0.15), Inches(1.15), Inches(6.2), Inches(5.8))

# Right half — text (6.6 wide)
txt(s, "مهارات الوكيل",
    Inches(6.7), Inches(0.15), Inches(6.4), Inches(0.6),
    size=28, bold=True, color=LIME)
txt(s, "(Agent Skills for Power BI)",
    Inches(6.7), Inches(0.73), Inches(6.4), Inches(0.4),
    size=15, color=CYAN, rtl=False, align=PP_ALIGN.RIGHT)
rect(s, Inches(6.7), Inches(1.2), Inches(6.1), Inches(0.04), LIME)

steps = [
    ("①", "مدخل طبيعي  (Natural Input)",
     "نص أو لقطة شاشة  (Text or Screenshot)", LIME),
    ("②", "الوكيل الذكي  (AI Agent)",
     "يفهم المطلوب ويخطط", CYAN),
    ("③", "موديل الأعمال  (Semantic Model)",
     "يبني الجداول والعلاقات والداكس", LIME),
    ("④", "بناء التقرير  (Report Building)",
     "صفحات ومرئيات وKPIs", CYAN),
    ("⑤", "نشر في فابريك  (Publish to Fabric)",
     "جاهز للاستخدام فوراً", LIME),
]
for i, (n, title, desc, col) in enumerate(steps):
    top = Inches(1.35 + i * 1.18)
    rect(s, Inches(6.7), top, Inches(6.3), Inches(1.05), BG)
    rect(s, Inches(6.7), top, Inches(0.06), Inches(1.05), col)

    txt(s, n, Inches(6.8), top + Inches(0.05), Inches(0.5), Inches(0.45),
        size=18, bold=True, color=col, align=PP_ALIGN.LEFT, rtl=False)
    txt(s, title, Inches(7.3), top + Inches(0.05), Inches(5.5), Inches(0.45),
        size=14, bold=True, color=WHITE)
    txt(s, desc, Inches(7.3), top + Inches(0.5), Inches(5.5), Inches(0.45),
        size=12, color=DIMMED)

logo(s); num(s, 4)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — FABRIC APPS: BEFORE/AFTER + visual comparison
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), CYAN)

txt(s, "تطبيقات فابريك  (Fabric Apps for Semantic Models)",
    Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.65),
    size=26, bold=True, color=CYAN)
txt(s, "من موديل بيانات إلى تطبيق ويب كامل — بكلام طبيعي",
    Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.42),
    size=17, color=WHITE)

# BEFORE column
rect(s, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.85), RED_DIM)
rect(s, Inches(0.4), Inches(1.35), Inches(5.9), Inches(0.07), RGBColor(0xCC, 0x22, 0x22))
txt(s, "قبل ❌  (Before Fabric Apps)",
    Inches(0.5), Inches(1.45), Inches(5.7), Inches(0.55),
    size=18, bold=True, color=RGBColor(0xFF, 0x66, 0x66), align=PP_ALIGN.CENTER)

before_items = [
    ("👨‍💻", "مطور واجهة أمامية\n(Frontend Developer)", "أسابيع من التصميم"),
    ("⚙️", "مطور خوادم\n(Backend Developer)", "API وقواعد البيانات"),
    ("🎨", "ديزاينر\n(UX Designer)", "تجربة المستخدم"),
    ("🔗", "واجهة برمجية\n(API Integration)", "ربط كل شي ببعض"),
    ("⏰", "3-6 أسابيع من الشغل", "فريق كامل مطلوب"),
]
for i, (icon, title, sub) in enumerate(before_items):
    top = Inches(2.1 + i * 0.95)
    txt(s, icon, Inches(0.55), top, Inches(0.5), Inches(0.6),
        size=20, align=PP_ALIGN.LEFT, rtl=False)
    txt(s, title, Inches(1.15), top, Inches(4.9), Inches(0.45),
        size=13, bold=True, color=RGBColor(0xFF, 0xAA, 0xAA))
    txt(s, sub, Inches(1.15), top + Inches(0.43), Inches(4.9), Inches(0.38),
        size=11, color=DIMMED)

# Center arrow
txt(s, "→", Inches(6.4), Inches(4.0), Inches(0.5), Inches(0.8),
    size=36, bold=True, color=LIME, align=PP_ALIGN.CENTER, rtl=False)

# AFTER column
rect(s, Inches(7.0), Inches(1.35), Inches(5.9), Inches(5.85), GRN_DIM)
rect(s, Inches(7.0), Inches(1.35), Inches(5.9), Inches(0.07), LIME)
txt(s, "بعد ✅  (With Fabric Apps)",
    Inches(7.1), Inches(1.45), Inches(5.7), Inches(0.55),
    size=18, bold=True, color=LIME, align=PP_ALIGN.CENTER)

after_items = [
    ("🤖", "وكيل ذكي واحد\n(One AI Agent)", "يفهم طلبك بالكلام"),
    ("🗄️", "موديل أعمال جاهز\n(Semantic Model)", "يستخدمه مباشرة"),
    ("💬", "طلب طبيعي\n(Natural Language)", "تصف ما تريد فقط"),
    ("🚀", "تطبيق ويب كامل\n(Fabric-native App)", "جاهز ومنشور"),
    ("⚡", "أيام وليس أسابيع", "بدون فريق تطوير"),
]
for i, (icon, title, sub) in enumerate(after_items):
    top = Inches(2.1 + i * 0.95)
    txt(s, icon, Inches(7.15), top, Inches(0.5), Inches(0.6),
        size=20, align=PP_ALIGN.LEFT, rtl=False)
    txt(s, title, Inches(7.75), top, Inches(4.9), Inches(0.45),
        size=13, bold=True, color=RGBColor(0xAA, 0xFF, 0x88))
    txt(s, sub, Inches(7.75), top + Inches(0.43), Inches(4.9), Inches(0.38),
        size=11, color=DIMMED)

logo(s); num(s, 5)


# ═══════════════════════════════════════════════════════
# SLIDE 6 — MCP: Left image panel | Right two server cards
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), LIME)

# LEFT PANEL — image at its natural aspect ratio (square-ish)
# 6" wide will give ~6" tall for a square image — fits nicely in left half
rect(s, Inches(0), Inches(0.06), Inches(6.55), Inches(7.38), CARD)
txt(s, "بروتوكول السياق النموذجي",
    Inches(0.2), Inches(0.15), Inches(6.2), Inches(0.55),
    size=20, bold=True, color=LIME, align=PP_ALIGN.RIGHT)
txt(s, "(MCP — Model Context Protocol)",
    Inches(0.2), Inches(0.68), Inches(6.2), Inches(0.4),
    size=13, color=CYAN, align=PP_ALIGN.RIGHT, rtl=False)
# Image: 6" wide, auto height preserves aspect — stays in left panel
img(s, "mcp_architecture_diagram.png",
    Inches(0.25), Inches(1.15), Inches(6.1))  # width only — no forced height

# USB-C caption at bottom of left panel
rect(s, Inches(0.2), Inches(6.85), Inches(6.15), Inches(0.45), BG)
txt(s, "زي USB-C للذكاء الاصطناعي — معيار موحّد  (Universal AI Protocol)",
    Inches(0.25), Inches(6.88), Inches(6.1), Inches(0.4),
    size=11, color=DIMMED, align=PP_ALIGN.CENTER)

# RIGHT PANEL — title + two server explanation cards
rect(s, Inches(6.65), Inches(0.06), Inches(6.68), Inches(7.38), BG)
txt(s, "نوعان من الخوادم",
    Inches(6.75), Inches(0.15), Inches(6.4), Inches(0.55),
    size=22, bold=True, color=LIME, align=PP_ALIGN.LEFT, rtl=False)
txt(s, "(Two Types of MCP Servers)",
    Inches(6.75), Inches(0.68), Inches(6.4), Inches(0.4),
    size=14, color=CYAN, align=PP_ALIGN.LEFT, rtl=False)
rect(s, Inches(6.75), Inches(1.15), Inches(6.2), Inches(0.04), LIME)

server_cards = [
    ("خادم التحليل البعيد ☁️",
     "(Remote Power BI MCP Server)",
     ["◆ يتصل بموديلات الأعمال المنشورة  (Published Semantic Models)",
      "◆ يولّد استعلامات داكس تلقائياً  (Auto DAX Queries)",
      "◆ يجاوب أسئلة البيانات بالكلام الطبيعي",
      "◆ يحترم صلاحيات الأمان  (Row-Level Security)",
      "◆ لا تحتاج وصول مباشر للبيانات"],
     CYAN),
    ("خادم التطوير المحلي 💻",
     "(Local Power BI Modeling MCP Server)",
     ["◆ يشتغل على ملفات المشروع المحلية  (.pbip files)",
      "◆ عمليات جماعية  (Bulk Operations) — تغيير أعمدة وعلاقات دفعة واحدة",
      "◆ إنشاء مقاييس داكس  (DAX Measures) تلقائياً بالكلام",
      "◆ يعمل مع بيئات التطوير  (VS Code / Claude Desktop)",
      "◆ وصول كامل لمشروع Power BI"],
     LIME),
]
for i, (title, sub, bullets, col) in enumerate(server_cards):
    top = Inches(1.28 + i * 3.05)
    rect(s, Inches(6.75), top, Inches(6.3), Inches(2.8), CARD)
    rect(s, Inches(6.75), top, Inches(0.06), Inches(2.8), col)
    txt(s, title, Inches(6.95), top + Inches(0.08),
        Inches(5.95), Inches(0.52), size=16, bold=True, color=col)
    txt(s, sub, Inches(6.95), top + Inches(0.58),
        Inches(5.95), Inches(0.38), size=11, color=DIMMED, rtl=False, align=PP_ALIGN.LEFT)
    rect(s, Inches(6.9), top + Inches(0.95), Inches(5.9), Inches(0.03), col)
    for j, b in enumerate(bullets):
        txt(s, b, Inches(6.9), top + Inches(1.05 + j * 0.37),
            Inches(6.1), Inches(0.35), size=11, color=WHITE)

logo(s); num(s, 6)
for title, sub, bullets, col, l in server_cards:
    rect(s, l, Inches(4.9), Inches(6.15), Inches(2.4), CARD)
    rect(s, l, Inches(4.9), Inches(0.06), Inches(2.4), col)
    txt(s, title, l + Inches(0.15), Inches(4.95), Inches(5.85), Inches(0.48),
        size=16, bold=True, color=col)
    txt(s, sub, l + Inches(0.15), Inches(5.4), Inches(5.85), Inches(0.35),
        size=11, color=DIMMED, rtl=False, align=PP_ALIGN.RIGHT)
    for j, b in enumerate(bullets):
        txt(s, b, l + Inches(0.15), Inches(5.75 + j * 0.37),
            Inches(5.85), Inches(0.35), size=11, color=WHITE)

logo(s); num(s, 6)


# ══════════════════════════════════════════════════════════
# SLIDE 7 — FABRIC IQ: RIGHT IMAGE | LEFT TEXT
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), CYAN)

# Left text side (6.2 wide)
txt(s, "ذكاء فابريك  (Fabric IQ)",
    Inches(0.4), Inches(0.15), Inches(6.0), Inches(0.65),
    size=28, bold=True, color=CYAN)
txt(s, "طبقة الذكاء الدلالي للمنصة",
    Inches(0.4), Inches(0.78), Inches(6.0), Inches(0.42),
    size=17, color=WHITE)
rect(s, Inches(0.4), Inches(1.28), Inches(5.7), Inches(0.04), CYAN)

points = [
    (LIME, "المشكلة",
     "الذكاء الاصطناعي بدون سياق يخمّن\nعلى أساس أسماء الجداول — وليس منطق الأعمال الحقيقي"),
    (CYAN, "الحل",
     "ذكاء فابريك يوحّد المقاييس والعلاقات\nوتعريفات الأعمال في طبقة واحدة مشتركة  (Unified Semantic Layer)"),
    (LIME, "النتيجة",
     "كل الوكلاء يشتغلون على نفس فهم الأعمال\nبدون تناقض أو أخطاء  (No Contradictions)"),
]
for i, (col, label, desc) in enumerate(points):
    top = Inches(1.4 + i * 1.85)
    rect(s, Inches(0.4), top, Inches(6.0), Inches(1.68), CARD)
    rect(s, Inches(0.4), top, Inches(0.06), Inches(1.68), col)
    txt(s, label, Inches(0.6), top + Inches(0.1), Inches(5.6), Inches(0.45),
        size=15, bold=True, color=col)
    txt(s, desc, Inches(0.6), top + Inches(0.55), Inches(5.6), Inches(1.1),
        size=13, color=WHITE)

# Bottom IQ ecosystem row
rect(s, Inches(0.4), Inches(6.85), Inches(6.0), Inches(0.5), BG)
for j, (label, col) in enumerate([
    ("Work IQ  (M365)", CYAN),
    ("Fabric IQ  (Data)", LIME),
    ("Foundry IQ  (Docs)", CYAN),
]):
    lx = Inches(0.5 + j * 2.0)
    rect(s, lx, Inches(6.85), Inches(1.85), Inches(0.45), CARD, line_color=col)
    txt(s, label, lx, Inches(6.87), Inches(1.85), Inches(0.42),
        size=11, bold=True, color=col, align=PP_ALIGN.CENTER, rtl=False)

# Right side — image
rect(s, Inches(6.6), Inches(0.06), Inches(6.73), Inches(7.38), CARD)
txt(s, "منظومة مايكروسوفت آي كيو",
    Inches(6.7), Inches(0.15), Inches(6.5), Inches(0.55),
    size=17, bold=True, color=LIME, align=PP_ALIGN.LEFT, rtl=False)
txt(s, "(Microsoft IQ Ecosystem)",
    Inches(6.7), Inches(0.68), Inches(6.5), Inches(0.38),
    size=13, color=CYAN, align=PP_ALIGN.LEFT, rtl=False)
img(s, "fabric_iq_ecosystem.png",
    Inches(6.65), Inches(1.15), Inches(6.5), Inches(6.15))

logo(s); num(s, 7)


# ══════════════════════════════════════════════════════════
# SLIDE 8 — IMPACT: 4 Quadrant Cards with Stats
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), LIME)

txt(s, "ماذا يعني هذا التطور عليك؟  (What Does This Mean for You?)",
    Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.65),
    size=26, bold=True, color=LIME)

quadrants = [
    ("80%", "تخفيض في الأعمال المتكررة\n(Reduction in Repetitive Work)",
     "إنجازات 2026 من المؤسسات اللي طبّقت الأجنتس", LIME,
     Inches(0.4), Inches(1.0)),
    ("∞", "موديل الأعمال هو المهارة الأهم\n(Semantic Model is King)",
     "الوكيل يكتب الداكس — لكن مين يبني الموديل الصح؟", CYAN,
     Inches(6.85), Inches(1.0)),
    ("↑", "دور المحلل يتطور مش يختفي\n(The Analyst Role Evolves)",
     "من تنفيذ المهام → الإشراف والجودة والاستراتيجية", LIME,
     Inches(0.4), Inches(4.35)),
    ("🔒", "الحوكمة صارت ميزة تنافسية\n(Governance = Competitive Edge)",
     "المؤسسات ذات الحوكمة الجيدة ستكسب السباق", CYAN,
     Inches(6.85), Inches(4.35)),
]

for stat, title, desc, col, l, t in quadrants:
    rect(s, l, t, Inches(6.1), Inches(3.1), CARD)
    rect(s, l, t, Inches(6.1), Inches(0.07), col)

    # Big stat
    txt(s, stat, l, t + Inches(0.1), Inches(2.2), Inches(1.2),
        size=54, bold=True, color=col, align=PP_ALIGN.LEFT, rtl=False)

    # Title
    txt(s, title, l + Inches(2.2), t + Inches(0.15), Inches(3.8), Inches(1.1),
        size=15, bold=True, color=WHITE)

    # Divider
    rect(s, l + Inches(0.2), t + Inches(1.35), Inches(5.7), Inches(0.03), col)

    # Description
    txt(s, desc, l + Inches(0.2), t + Inches(1.5), Inches(5.7), Inches(1.45),
        size=13, color=DIMMED)

logo(s); num(s, 8)


# ══════════════════════════════════════════════════════════
# SLIDE 9 — SUMMARY RECAP: Full visual checklist
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, Inches(0), CYAN)
bar(s, Inches(7.44), LIME)

txt(s, "خلاصة الكلام  (Summary)",
    Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.65),
    size=30, bold=True, color=WHITE)

recaps = [
    (LIME, "01", "مهارات الوكيل  (Agent Skills for Power BI)",
     "بناء موديل أعمال كامل وتقارير بالكلام الطبيعي أو لقطة شاشة",
     "موديل أعمال + داكس + تقرير + نشر = وكيل واحد"),
    (CYAN, "02", "تطبيقات فابريك  (Fabric Apps for Semantic Models)",
     "تطبيق ويب أصيل فوق موديل البيانات — بدون فريق تطوير",
     "أيام بدل أسابيع — وكيل ذكي واحد يبني كل شي"),
    (LIME, "03", "بروتوكول السياق النموذجي  (MCP Servers)",
     "جسر موحّد بين الوكلاء وأنظمة باور بي آي وفابريك",
     "خادم بعيد للتحليل + خادم محلي للتطوير"),
    (CYAN, "04", "ذكاء فابريك  (Fabric IQ)",
     "الطبقة الذكية الأساسية — كل الوكلاء يفهمون الأعمال بنفس الطريقة",
     "Work IQ + Fabric IQ + Foundry IQ = Microsoft IQ"),
]

for i, (col, n, title, desc, detail) in enumerate(recaps):
    top = Inches(0.95 + i * 1.55)
    rect(s, Inches(0.4), top, Inches(12.5), Inches(1.4), CARD)
    rect(s, Inches(0.4), top, Inches(0.07), Inches(1.4), col)

    # Number
    rect(s, Inches(0.6), top + Inches(0.2), Inches(0.65), Inches(0.65), col)
    txt(s, n, Inches(0.6), top + Inches(0.2), Inches(0.65), Inches(0.65),
        size=17, bold=True, color=BG, align=PP_ALIGN.CENTER, rtl=False)

    # Title
    txt(s, title, Inches(1.45), top + Inches(0.08), Inches(8.5), Inches(0.55),
        size=16, bold=True, color=col)
    txt(s, desc, Inches(1.45), top + Inches(0.62), Inches(8.5), Inches(0.45),
        size=12, color=WHITE)

    # Right detail tag
    rect(s, Inches(10.1), top + Inches(0.3), Inches(2.65), Inches(0.65), BG,
         line_color=col)
    txt(s, detail, Inches(10.15), top + Inches(0.33), Inches(2.55), Inches(0.58),
        size=10, color=col, align=PP_ALIGN.LEFT, rtl=False)

logo(s); num(s, 9)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — CTA
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, Inches(0), Inches(0.06), Inches(0.08), Inches(7.38), CYAN)
rect(s, Inches(13.25), Inches(0.06), Inches(0.08), Inches(7.38), LIME)
bar(s, Inches(0), LIME)
bar(s, Inches(7.44), CYAN)

txt(s, "عصر الذكاء الاصطناعي الوكيل",
    Inches(1.0), Inches(0.9), Inches(11.3), Inches(1.0),
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "(Agentic AI Era)",
    Inches(1.0), Inches(1.85), Inches(11.3), Inches(0.6),
    size=22, color=LIME, align=PP_ALIGN.CENTER, rtl=False)

rect(s, Inches(1.5), Inches(2.55), Inches(10.3), Inches(0.04), CYAN)

txt(s, "ما جاء ليحل محلك — جاء ليضاعف قدرتك 🚀",
    Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.75),
    size=24, bold=True, color=LIME, align=PP_ALIGN.CENTER)

# 4 final points
cta_pts = [
    ("✅", "مهارات الوكيل\n(Agent Skills)", LIME),
    ("✅", "تطبيقات فابريك\n(Fabric Apps)", CYAN),
    ("✅", "بروتوكول MCP\n(MCP Protocol)", LIME),
    ("✅", "ذكاء فابريك\n(Fabric IQ)", CYAN),
]
for i, (icon, label, col) in enumerate(cta_pts):
    l = Inches(0.8 + i * 3.1)
    rect(s, l, Inches(3.65), Inches(2.9), Inches(1.7), CARD)
    rect(s, l, Inches(3.65), Inches(2.9), Inches(0.06), col)
    txt(s, icon + "\n" + label, l, Inches(3.7), Inches(2.9), Inches(1.6),
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER, rtl=False)

rect(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.04), LIME)

txt(s, "اشترك في القناة وخلّي اللايك ❤️",
    Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.65),
    size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "@HakamDataStudio",
    Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.75),
    size=32, bold=True, color=CYAN, align=PP_ALIGN.CENTER, rtl=False)

logo(s); num(s, 10)


# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
out = SCRIPT_DIR / "agentic-ai-presentation.pptx"
prs.save(str(out))
print(f"[OK] Saved: {out}")
print(f"[OK] Slides: {len(prs.slides)}")
